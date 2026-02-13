from pptx import Presentation


def read_pptx(path: str) -> list[dict]:
    prs = Presentation(path)
    pages = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        text = "\n".join(texts)
        if text.strip():
            pages.append({"text": text, "page": i})
    return pages
