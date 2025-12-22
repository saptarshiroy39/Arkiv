# Arkiv Extractor
# 1. Extract text from files
# 2. Process images
# 3. Handle errors
# 4. Sanitize text
# 5. Return text

import csv
import io
import os
from typing import Callable, Optional

import pdfplumber
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

from .config import (
    logger,
    PDF_EXTENSIONS,
    IMAGE_EXTENSIONS,
    WORD_EXTENSIONS,
    EXCEL_EXTENSIONS,
    POWERPOINT_EXTENSIONS,
    TEXT_EXTENSIONS,
    MARKDOWN_EXTENSIONS,
    CSV_EXTENSIONS,
)
from .processor import sanitize_text, process_image, simplify_error_message


def _handle_error(func_name: str, filename: str, error: Exception) -> str:
    """Centralized error formatting."""
    error_msg = simplify_error_message(error)
    logger.error(f"{func_name} failed for '{filename}': {str(error)}")
    return f"[{func_name}: {filename} - Failed to extract text: {error_msg}]"


def _process_pdf(file_bytes: bytes, filename: str, **kwargs) -> str:
    """Extract text from PDF using pdfplumber."""
    try:
        text_parts = [f"=== PDF Document: {filename} ==="]
        extracted_any = False
        
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_parts.append(page_text)
                        extracted_any = True
                except Exception as e:
                    logger.warning(f"Page extraction error in {filename}: {e}")
                    continue
                    
        if not extracted_any:
            logger.warning(f"PDF {filename} contained no extractable text (scanned?)")
            return f"[PDF: {filename} - No text extraction likely due to scanned image/protection]"
            
        logger.info(f"Successfully processed PDF: {filename}")
        return sanitize_text("\n".join(text_parts))
        
    except Exception as e:
        return _handle_error("PDF Processing", filename, e)


def _process_word(file_bytes: bytes, filename: str, **kwargs) -> str:
    """Extract text from Word (.docx) files."""
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
        texts = [f"=== Word Document: {filename} ==="]
        
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
                
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text for cell in row.cells)
                if row_text.strip():
                    texts.append(row_text)
                    
        return sanitize_text("\n".join(texts))
        
    except Exception as e:
        return _handle_error("Word Processing", filename, e)


def _process_excel(file_bytes: bytes, filename: str, **kwargs) -> str:
    """Extract text from Excel (.xlsx) files."""
    try:
        wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
        texts = [f"=== Excel Spreadsheet: {filename} ==="]
        
        for sheet_name in wb.sheetnames:
            texts.append(f"\n--- Sheet: {sheet_name} ---")
            sheet = wb[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                # Convert None to "" and other types to string
                row_values = [str(cell) if cell is not None else "" for cell in row]
                # Only add row if it has content
                if any(val.strip() for val in row_values):
                    texts.append("\t".join(row_values))
                    
        return sanitize_text("\n".join(texts))
        
    except Exception as e:
        return _handle_error("Excel Processing", filename, e)


def _process_powerpoint(file_bytes: bytes, filename: str, **kwargs) -> str:
    """Extract text from PowerPoint (.pptx) files."""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        texts = [f"=== PowerPoint: {filename} ==="]
        
        for i, slide in enumerate(prs.slides, 1):
            texts.append(f"\n--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
                    
        return sanitize_text("\n".join(texts))
        
    except Exception as e:
        return _handle_error("PowerPoint Processing", filename, e)


def _process_csv(file_bytes: bytes, filename: str, **kwargs) -> str:
    """Extract text from CSV files."""
    try:
        text_content = file_bytes.decode('utf-8', errors='ignore')
        texts = [f"=== CSV File: {filename} ==="]
        
        reader = csv.reader(text_content.splitlines())
        for row in reader:
            if any(cell.strip() for cell in row):
                texts.append("\t".join(row))
                
        return sanitize_text("\n".join(texts))
        
    except Exception as e:
        return _handle_error("CSV Processing", filename, e)


def _process_text_or_markdown(file_bytes: bytes, filename: str, **kwargs) -> str:
    """Extract text from generic Text or Markdown files."""
    try:
        content = file_bytes.decode('utf-8', errors='ignore')
        return sanitize_text(f"=== Document: {filename} ===\n" + content)
    except Exception as e:
        return _handle_error("Text Processing", filename, e)


def _process_image(file_bytes: bytes, filename: str, api_key: str = None, **kwargs) -> str:
    """Process images using AI Vision (Gemini)."""
    try:
        description = process_image(file_bytes, filename, api_key=api_key)
        if description:
            return description
        return ""
    except Exception as e:
        return _handle_error("Image Processing", filename, e)



HandlerFunction = Callable[..., str]

_EXTENSION_HANDLERS: dict[tuple[str, ...], HandlerFunction] = {
    PDF_EXTENSIONS:        _process_pdf,
    WORD_EXTENSIONS:       _process_word,
    EXCEL_EXTENSIONS:      _process_excel,
    POWERPOINT_EXTENSIONS: _process_powerpoint,
    IMAGE_EXTENSIONS:      _process_image,
    CSV_EXTENSIONS:        _process_csv,
    TEXT_EXTENSIONS:       _process_text_or_markdown,
    MARKDOWN_EXTENSIONS:   _process_text_or_markdown,
}

def extract_text_from_file(file_bytes: bytes, filename: str, api_key: str = None) -> str:
    """
    Main Entry Point: Dispatches file to the correct handler based on extension.
    """
    filename_lower = filename.lower()
    _, ext = os.path.splitext(filename_lower)
    
    
    for extension_group, handler in _EXTENSION_HANDLERS.items():
        if ext in extension_group:
            return handler(file_bytes, filename, api_key=api_key)
            
    raise ValueError(f"Unsupported file type: {ext}")
