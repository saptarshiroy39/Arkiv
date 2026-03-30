from pptx import Presentation


def read_pptx(path: str) -> list[dict]:
    prs = Presentation(path)
    pages = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        text = "\n".join(texts)
        if text.strip():
            pages.append({"text": text})
    return pages
